import os
import json
from typing import Iterator
from abc import ABC

import pandas as pd
from langchain_core.documents import Document
from langchain_community.document_loaders.base import BaseBlobParser
from langchain_community.document_loaders.blob_loaders import Blob


class TableParser(BaseBlobParser, ABC):
    """Parsa CSV/XLSX. Per XLSX emette 1 Document per foglio."""

    def _get_source(self, blob: Blob) -> str:
        p = getattr(blob, "path", None) or getattr(blob, "source", None) or ""
        try:
            return os.path.basename(p) if p else ""
        except Exception:
            return str(p)

    def lazy_parse(self, blob: Blob) -> Iterator[Document]:
        with blob.as_bytes_io() as file:
            if blob.mimetype == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                # Un Document per sheet
                df_dict = pd.read_excel(file, sheet_name=None)
                for sheet_name, sheet_df in df_dict.items():
                    if sheet_df is None or sheet_df.empty:
                        continue
                    records = sheet_df.to_dict("records")
                    # opzionale: traccia la provenienza riga->foglio
                    for r in records:
                        r.setdefault("_sheet", sheet_name)

                    yield Document(
                        page_content=json.dumps(records, ensure_ascii=False),
                        metadata={
                            "source": self._get_source(blob),
                            "mimetype": blob.mimetype,
                            "sheet_name": sheet_name,
                            "row_count": len(records),
                            "parser": "TableParser",
                        },
                    )

            elif blob.mimetype == "text/csv":
                df = pd.read_csv(file)
                records = df.to_dict("records") if not df.empty else []
                for r in records:
                    r.setdefault("_sheet", "CSV")

                yield Document(
                    page_content=json.dumps(records, ensure_ascii=False),
                    metadata={
                        "source": self._get_source(blob),
                        "mimetype": blob.mimetype,
                        "sheet_name": "CSV",
                        "row_count": len(records),
                        "parser": "TableParser",
                    },
                )
            else:
                raise ValueError(f"Unsupported table mime type: {blob.mimetype}")


class PowerPointParser(BaseBlobParser, ABC):
    """Estrae il testo dalle slide e include metadati utili."""

    def _get_source(self, blob: Blob) -> str:
        p = getattr(blob, "path", None) or getattr(blob, "source", None) or ""
        try:
            return os.path.basename(p) if p else ""
        except Exception:
            return str(p)

    def lazy_parse(self, blob: Blob) -> Iterator[Document]:
        pptx_mime_types = [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # .pptx
            "application/vnd.ms-powerpoint",  # .ppt
            "application/powerpoint",  # Alternative .ppt
        ]
        if blob.mimetype not in pptx_mime_types:
            raise ValueError(f"Unsupported mime type: {blob.mimetype}")

        with blob.as_bytes_io() as file_obj:
            import pptx as _pptx

            presentation = _pptx.Presentation(file_obj)

            all_text = []
            slide_contents = {}

            for i, slide in enumerate(presentation.slides, 1):
                slide_text = []
                title = ""

                # Title (se presente)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and getattr(shape, "is_title", False):
                        title = shape.text or ""
                        break

                # Tutti i testi della slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)

                slide_content = "\n".join(slide_text)
                all_text.append(slide_content)
                slide_contents[f"Slide {i}"] = {"title": title, "content": slide_content}

            full_text = "\n\n".join(all_text)

            yield Document(
                page_content=full_text,
                metadata={
                    "source": self._get_source(blob),
                    "mimetype": blob.mimetype,
                    "parser": "PowerPointParser",
                },
            )


# class EmailParser(BaseBlobParser, ABC):
#     """Parsa .eml/.msg mantenendo header essenziali, body pulito e info allegati."""

#     def _get_source(self, blob: Blob) -> str:
#         p = getattr(blob, "path", None) or getattr(blob, "source", None) or ""
#         try:
#             return os.path.basename(p) if p else ""
#         except Exception:
#             return str(p)

#     def _html_to_text(self, html: str) -> str:
#         if not html:
#             return ""
#         try:
#             from bs4 import BeautifulSoup
#             return BeautifulSoup(html, "html.parser").get_text(separator="\n")
#         except Exception:
#             # Fallback leggero senza bs4
#             import re, html as _html
#             text = re.sub(r"(?is)<(script|style).*?>.*?(</\1>)", "", html)
#             text = re.sub(r"(?s)<br\s*/?>", "\n", text)
#             text = re.sub(r"(?s)</p\s*>", "\n\n", text)
#             text = re.sub(r"(?s)<.*?>", "", text)
#             return _html.unescape(text)

#     def _normalize(self, s: str) -> str:
#         return (s or "").replace("\r\n", "\n").replace("\r", "\n").strip()

#     def lazy_parse(self, blob: Blob) -> Iterator[Document]:
#         email_mime_types = [
#             "message/rfc822",                 # .eml
#             "application/vnd.ms-outlook",     # .msg
#             "application/octet-stream",       # talvolta .eml/.msg
#             "application/x-ole-storage",      # alcuni .msg
#         ]

#         path_lower = ((blob.path or getattr(blob, "source", "")) or "").lower()
#         is_eml = (blob.mimetype == "message/rfc822") or path_lower.endswith(".eml")
#         is_msg = (blob.mimetype in ("application/vnd.ms-outlook", "application/x-ole-storage")) or path_lower.endswith(".msg")

#         if not (blob.mimetype in email_mime_types or is_eml or is_msg):
#             raise ValueError(f"Unsupported mime type: {blob.mimetype}")

#         with blob.as_bytes_io() as file_obj:
#             sender = recipients = cc = bcc = subject = date_str = ""
#             body_text = ""
#             attachments_meta = []

#             if is_eml:
#                 import email
#                 from email import policy

#                 msg = email.message_from_binary_file(file_obj, policy=policy.default)
#                 subject = self._normalize(msg.get("Subject", ""))
#                 sender = self._normalize(msg.get("From", ""))
#                 recipients = self._normalize(msg.get("To", ""))
#                 cc = self._normalize(msg.get("Cc", ""))
#                 bcc = self._normalize(msg.get("Bcc", ""))
#                 date_str = self._normalize(msg.get("Date", ""))

#                 if msg.is_multipart():
#                     # preferisci text/plain; se assente, ripulisci l'HTML
#                     for part in msg.iter_parts():
#                         disp = (part.get("Content-Disposition") or "").lower()
#                         ctype = (part.get_content_type() or "").lower()
#                         if "attachment" in disp:
#                             fname = part.get_filename() or ""
#                             try:
#                                 payload = part.get_content()
#                                 size = len(payload.encode("utf-8")) if isinstance(payload, str) else len(payload or b"")
#                             except Exception:
#                                 size = None
#                             attachments_meta.append({"filename": fname, "content_type": ctype, "size": size})
#                             continue
#                         if ctype == "text/plain" and not body_text:
#                             body_text = self._normalize(part.get_content())
#                         elif ctype == "text/html" and not body_text:
#                             body_text = self._normalize(self._html_to_text(part.get_content()))
#                 else:
#                     ctype = (msg.get_content_type() or "").lower()
#                     if ctype == "text/plain":
#                         body_text = self._normalize(msg.get_content())
#                     elif ctype == "text/html":
#                         body_text = self._normalize(self._html_to_text(msg.get_content()))

#             elif is_msg:
#                 # Outlook .msg via extract_msg richiede un path temporaneo
#                 import tempfile
#                 import os as _os
#                 try:
#                     import extract_msg
#                 except ImportError as e:
#                     raise ImportError("Per leggere .msg serve il pacchetto 'extract_msg'.") from e

#                 file_obj.seek(0)
#                 with tempfile.NamedTemporaryFile(delete=False) as temp:
#                     temp.write(file_obj.read())
#                     temp_path = temp.name

#                 try:
#                     msg = extract_msg.Message(temp_path)
#                     msg.populate()  # assicura che i campi siano caricati

#                     subject = self._normalize(getattr(msg, "subject", "") or "")
#                     sender = self._normalize(getattr(msg, "sender", "") or getattr(msg, "sender", ""))
#                     recipients = self._normalize(getattr(msg, "to", "") or "")
#                     cc = self._normalize(getattr(msg, "cc", "") or "")
#                     bcc = self._normalize(getattr(msg, "bcc", "") or "")
#                     date_obj = getattr(msg, "date", None)
#                     date_str = self._normalize(str(date_obj)) if date_obj else ""

#                     # Ordine di preferenza: body (plain) -> htmlBody -> rtfBody
#                     plain = getattr(msg, "body", "") or ""
#                     html = getattr(msg, "htmlBody", "") or ""
#                     rtf = getattr(msg, "rtfBody", "") or ""

#                     if plain.strip():
#                         body_text = self._normalize(plain)
#                     elif html.strip():
#                         body_text = self._normalize(self._html_to_text(html))
#                     elif rtf.strip():
#                         try:
#                             from striprtf.striprtf import rtf_to_text
#                             body_text = self._normalize(rtf_to_text(rtf))
#                         except Exception:
#                             body_text = self._normalize(rtf)

#                     # Allegati
#                     atts = getattr(msg, "attachments", []) or []
#                     for att in atts:
#                         # extract_msg Attachment object
#                         fname = getattr(att, "longFilename", None) or getattr(att, "shortFilename", None) or ""
#                         try:
#                             data = getattr(att, "data", None)
#                             size = len(data) if data is not None else None
#                         except Exception:
#                             size = None
#                         attachments_meta.append({
#                             "filename": fname,
#                             "content_type": None,
#                             "size": size,
#                         })

#                     msg.close()
#                 finally:
#                     _os.unlink(temp_path)
#             else:
#                 raise ValueError(f"Unsupported email format: {blob.mimetype}")

#             # Componi il testo finale (header + body)
#             parts = []
#             if sender:
#                 parts.append(f"Da: {sender}")
#             if recipients:
#                 parts.append(f"A: {recipients}")
#             if cc:
#                 parts.append(f"CC: {cc}")
#             if bcc:
#                 parts.append(f"BCC: {bcc}")
#             if date_str:
#                 parts.append(f"Data: {date_str}")
#             if subject:
#                 parts.append(f"Oggetto: {subject}")
#             if body_text:
#                 parts.append("\n" + body_text)

#             # Riassunto allegati in coda (opzionale, utile all'indicizzazione)
#             if attachments_meta:
#                 parts.append("\n-- Allegati --")
#                 for a in attachments_meta:
#                     size_str = f" ({a['size']} bytes)" if a.get("size") else ""
#                     parts.append(f"- {a.get('filename','')}"+size_str)

#             full_content = "\n".join(parts)

#             yield Document(
#                 page_content=full_content,
#                 metadata={
#                     "source": self._get_source(blob),
#                     "mimetype": blob.mimetype,
#                     "parser": "EmailParser"
#                 },
#             )

